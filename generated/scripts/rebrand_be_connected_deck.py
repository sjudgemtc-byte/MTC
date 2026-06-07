from __future__ import annotations

import math
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


SOURCE_PPTX = Path("/Users/simonjudge/Downloads/be_connected_buying_selling_online_lesson_deck.pptx")
SOURCE_SLIDE_IMAGES = Path(
    "/Users/simonjudge/Documents/MTC/outputs/manual-20260606-branding/presentations/be-connected-branding/source_inspect/template-inspect/source-slides"
)
ARROW_IMAGE = Path(
    "/Users/simonjudge/Documents/MTC/outputs/manual-20260606-branding/presentations/be-connected-branding/template_media/image28.png"
)
OUTPUT_PPTX = Path(
    "/Users/simonjudge/Documents/MTC/04_Classroom_Activities/2026-06-06_MTC_Redfern_RED-01_Pathway_Be-Connected_Buying-Selling-Online_Presentation.pptx"
)

BRAND_COLORS = [
    RGBColor(255, 238, 72),   # yellow
    RGBColor(102, 226, 193),  # mint
    RGBColor(244, 109, 167),  # pink
    RGBColor(126, 203, 248),  # blue
    RGBColor(171, 138, 242),  # purple
]
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT_GREY = RGBColor(243, 243, 243)
MID_GREY = RGBColor(102, 102, 102)


def iter_source_content() -> list[dict[str, str]]:
    prs = Presentation(str(SOURCE_PPTX))
    slides: list[dict[str, str]] = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip().replace("\n", " ").strip()
                if text:
                    texts.append(text)
        title = texts[0] if texts else f"Slide {idx}"
        subtitle = texts[1] if len(texts) > 1 else ""
        slides.append(
            {
                "index": str(idx),
                "title": title,
                "subtitle": subtitle,
                "image": str(SOURCE_SLIDE_IMAGES / f"source-slide-{idx:02d}.png"),
            }
        )
    return slides


def set_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size, color, bold=False, font_name="Aptos"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_cover_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BLACK)

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(8.5), Inches(0.5), Inches(3.1), Inches(1.2))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BRAND_COLORS[0]
    accent.line.fill.background()

    accent2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(7.4), Inches(4.95), Inches(2.8), Inches(0.35))
    accent2.fill.solid()
    accent2.fill.fore_color.rgb = BRAND_COLORS[1]
    accent2.line.fill.background()

    if ARROW_IMAGE.exists():
        slide.shapes.add_picture(str(ARROW_IMAGE), Inches(8.15), Inches(1.05), width=Inches(2.2))

    title_box = add_textbox(slide, Inches(0.85), Inches(1.2), Inches(7.1), Inches(1.6), title, 24, WHITE, bold=True)
    title_box.text_frame.word_wrap = True

    sub_box = add_textbox(slide, Inches(0.9), Inches(3.1), Inches(6.6), Inches(1.0), subtitle, 14, WHITE)
    sub_box.text_frame.word_wrap = True

    footer = add_textbox(slide, Inches(0.9), Inches(6.45), Inches(5.8), Inches(0.35), "MTC FutureReady | Digital literacy classroom resource", 9, WHITE)
    footer.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    logo = add_textbox(slide, Inches(10.0), Inches(6.2), Inches(2.1), Inches(0.5), "mtc FutureReady", 18, WHITE, bold=True)
    logo.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT


def add_content_slide(prs: Presentation, slide_data: dict[str, str], color: RGBColor) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)

    header = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(0.45), Inches(2.9), Inches(0.6))
    header.fill.solid()
    header.fill.fore_color.rgb = color
    header.line.fill.background()

    title_box = add_textbox(slide, Inches(0.75), Inches(0.42), Inches(8.8), Inches(0.6), slide_data["title"], 20, BLACK, bold=True)
    title_box.text_frame.word_wrap = True

    if slide_data["subtitle"]:
        sub_box = add_textbox(slide, Inches(0.78), Inches(1.12), Inches(9.8), Inches(0.5), slide_data["subtitle"], 10.5, MID_GREY)
        sub_box.text_frame.word_wrap = True

    image_path = Path(slide_data["image"])
    if image_path.exists():
        pic_left = Inches(0.78)
        pic_top = Inches(1.7)
        max_width = Inches(11.0)
        max_height = Inches(4.9)
        pic = slide.shapes.add_picture(str(image_path), pic_left, pic_top)
        scale = min(max_width / pic.width, max_height / pic.height)
        pic.width = int(pic.width * scale)
        pic.height = int(pic.height * scale)
        pic.left = int(pic_left + (max_width - pic.width) / 2)
        pic.top = int(pic_top + (max_height - pic.height) / 2)

    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.75), Inches(6.68), Inches(10.2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = BLACK
    line.line.fill.background()

    footer = add_textbox(slide, Inches(0.8), Inches(6.75), Inches(4.2), Inches(0.28), "Be Connected | Buying and Selling Online", 8.5, MID_GREY)
    footer.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    slide_no = add_textbox(slide, Inches(10.55), Inches(6.68), Inches(0.45), Inches(0.28), slide_data["index"], 8.5, MID_GREY, bold=True)
    slide_no.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    if ARROW_IMAGE.exists():
        slide.shapes.add_picture(str(ARROW_IMAGE), Inches(11.1), Inches(6.46), width=Inches(0.34))


def build_deck() -> None:
    source = iter_source_content()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_cover_slide(prs, source[0]["title"], source[0]["subtitle"])
    for idx, slide_data in enumerate(source[1:], start=1):
        add_content_slide(prs, slide_data, BRAND_COLORS[(idx - 1) % len(BRAND_COLORS)])

    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PPTX))


if __name__ == "__main__":
    build_deck()
