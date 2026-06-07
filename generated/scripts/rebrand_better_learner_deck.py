from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SOURCE_PPTX = Path("/Users/simonjudge/Downloads/How_to_Be_a_Better_Learner_20_Minute_Lecture.pptx")
OUTPUT_PPTX = Path(
    "/Users/simonjudge/Documents/MTC/04_Classroom_Activities/2026-06-06_MTC_Redfern_RED-01_Pathway_How-to-Be-a-Better-Learner_Branded-Presentation.pptx"
)
BRAND_KIT = Path("/Users/simonjudge/Documents/MTC/07_Templates/Branding/MTC_Brand_Kit_2026-06-06")

LOGO_BLACK = BRAND_KIT / "MTC_Logo_Black_RGB.png"
LOGO_WHITE = BRAND_KIT / "MTC_Logo_White_RGB.png"

YELLOW = RGBColor(254, 241, 83)
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
CHARCOAL = RGBColor(63, 71, 79)
LIGHT_GREY = RGBColor(241, 243, 246)
MID_GREY = RGBColor(112, 112, 112)


def read_source_slides() -> list[dict[str, object]]:
    prs = Presentation(str(SOURCE_PPTX))
    slides: list[dict[str, object]] = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    texts.append(text)

        title = texts[1] if len(texts) > 1 else f"Slide {idx}"
        subtitle = texts[2] if idx == 1 and len(texts) > 2 else ""
        body_text = texts[3] if len(texts) > 3 else (texts[2] if len(texts) > 2 else "")
        note_text = texts[4] if len(texts) > 4 else ""

        bullets = []
        if body_text:
            for chunk in body_text.split("•"):
                cleaned = chunk.strip()
                if cleaned:
                    bullets.append(cleaned)

        slides.append(
            {
                "index": idx,
                "title": title,
                "subtitle": subtitle,
                "bullets": bullets,
                "note": note_text,
            }
        )
    return slides


def set_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_logo(slide, dark: bool = False) -> None:
    path = LOGO_WHITE if dark else LOGO_BLACK
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(0.45), Inches(6.55), width=Inches(1.9))


def add_page_number(slide, number: int, total: int, color: RGBColor) -> None:
    box = slide.shapes.add_textbox(Inches(12.2), Inches(0.28), Inches(0.65), Inches(0.3))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{number}/{total}"
    run.font.name = "Arial"
    run.font.size = Pt(12)
    run.font.color.rgb = color


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size,
    color,
    *,
    bold: bool = False,
    font_name: str = "Arial",
    align=PP_ALIGN.LEFT,
    vertical_anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = vertical_anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_note_strip(slide, text: str) -> None:
    strip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.65),
        Inches(6.38),
        Inches(12.0),
        Inches(0.72),
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = LIGHT_GREY
    strip.line.color.rgb = YELLOW
    strip.line.width = Pt(1.2)

    note_box = add_textbox(
        slide,
        Inches(1.0),
        Inches(6.53),
        Inches(11.2),
        Inches(0.34),
        text,
        9.5,
        CHARCOAL,
        font_name="Arial",
    )
    note_box.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT


def build_cover(prs: Presentation, slide_data: dict[str, object], total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, YELLOW)
    add_logo(slide, dark=False)
    add_page_number(slide, int(slide_data["index"]), total, CHARCOAL)

    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(5.95),
        Inches(13.333),
        Inches(1.55),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = BLACK
    band.line.fill.background()

    add_textbox(
        slide,
        Inches(0.8),
        Inches(0.75),
        Inches(10.8),
        Inches(1.1),
        str(slide_data["title"]),
        26,
        BLACK,
        bold=True,
        font_name="Arial",
    )

    add_textbox(
        slide,
        Inches(0.82),
        Inches(1.8),
        Inches(7.6),
        Inches(0.6),
        str(slide_data["subtitle"]),
        17,
        CHARCOAL,
        font_name="Arial",
    )

    bullet_box = slide.shapes.add_textbox(Inches(0.92), Inches(2.55), Inches(11.2), Inches(2.75))
    tf = bullet_box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, bullet in enumerate(slide_data["bullets"]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.name = "Arial"
        p.font.size = Pt(18)
        p.font.color.rgb = BLACK
        p.space_after = Pt(8)

    note = str(slide_data["note"])
    add_textbox(
        slide,
        Inches(0.85),
        Inches(6.22),
        Inches(11.6),
        Inches(0.55),
        note,
        11,
        WHITE,
        font_name="Arial",
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )


def build_content_slide(prs: Presentation, slide_data: dict[str, object], total: int, accent_yellow: bool) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = YELLOW if accent_yellow else WHITE
    set_background(slide, bg)
    add_logo(slide, dark=False)
    add_page_number(slide, int(slide_data["index"]), total, MID_GREY)

    title_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.56),
        Inches(0.52),
        Inches(0.25),
        Inches(0.88),
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = BLACK
    title_bar.line.fill.background()

    add_textbox(
        slide,
        Inches(0.92),
        Inches(0.46),
        Inches(11.1),
        Inches(0.95),
        str(slide_data["title"]),
        23,
        BLACK,
        bold=True,
        font_name="Arial",
    )

    bullet_box = slide.shapes.add_textbox(Inches(1.02), Inches(1.62), Inches(11.35), Inches(4.45))
    tf = bullet_box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, bullet in enumerate(slide_data["bullets"]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.name = "Arial"
        p.font.size = Pt(19)
        p.font.color.rgb = BLACK
        p.space_after = Pt(9)

    add_note_strip(slide, str(slide_data["note"]))


def build_closing_slide(prs: Presentation, slide_data: dict[str, object], total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BLACK)
    add_logo(slide, dark=True)
    add_page_number(slide, int(slide_data["index"]), total, WHITE)

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.78),
        Inches(0.72),
        Inches(11.55),
        Inches(5.55),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = YELLOW
    panel.line.fill.background()

    add_textbox(
        slide,
        Inches(1.15),
        Inches(1.02),
        Inches(9.8),
        Inches(0.9),
        str(slide_data["title"]),
        24,
        BLACK,
        bold=True,
        font_name="Arial",
    )

    bullet_box = slide.shapes.add_textbox(Inches(1.18), Inches(1.95), Inches(10.55), Inches(3.65))
    tf = bullet_box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, bullet in enumerate(slide_data["bullets"]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.name = "Arial"
        p.font.size = Pt(18)
        p.font.color.rgb = BLACK
        p.space_after = Pt(8)

    note_panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.82),
        Inches(6.05),
        Inches(11.5),
        Inches(0.58),
    )
    note_panel.fill.solid()
    note_panel.fill.fore_color.rgb = WHITE
    note_panel.line.fill.background()

    add_textbox(
        slide,
        Inches(1.05),
        Inches(6.17),
        Inches(11.0),
        Inches(0.24),
        str(slide_data["note"]),
        10,
        CHARCOAL,
        font_name="Arial",
    )


def build_deck() -> None:
    source = read_source_slides()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    total = len(source)
    build_cover(prs, source[0], total)
    for slide_data in source[1:-1]:
        accent = slide_data["index"] in {4, 8}
        build_content_slide(prs, slide_data, total, accent_yellow=accent)
    build_closing_slide(prs, source[-1], total)

    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PPTX))
    print(OUTPUT_PPTX)


if __name__ == "__main__":
    build_deck()
